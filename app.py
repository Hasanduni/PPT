import streamlit as st
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches
import requests
import os
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

# Retrieve the GlamaAI API key from environment variables
glamaai_api_key = os.getenv('GLAMAAI_API_KEY')
glamaai_api_url = 'https://glama.ai/api/gateway/openai/v1'  # Base URL for GlamaAI API

def generate_api_content(topic, description):
    """Fetch detailed content for the topic and description from GlamaAI API."""
    prompt = f"Generate a detailed explanation for the topic: {topic}\nDescription: {description}"

    headers = {
        'Authorization': f'Bearer {glamaai_api_key}',
        'Content-Type': 'application/json',
    }

    data = {
        'prompt': prompt,
        'max_tokens': 150,
        'temperature': 0.7
    }

    try:
        response = requests.post(glamaai_api_url + '/generate', json=data, headers=headers, timeout=30)
        response.raise_for_status()  # Raise HTTPError for bad responses
    except requests.exceptions.RequestException as e:
        st.error(f"Request failed: {e}")
        return "Content generation failed. Please try again."


    if response.status_code == 200:
        return response.json()['text']  # Assuming the response contains a 'text' field
    else:
        st.error("Error fetching content from GlamaAI API.")
        return "Content generation failed. Please try again."

def add_title_slide(prs, title, subtitle):
    """Add a title slide with a title and subtitle."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_chart_slide(prs, title):
    """Add a sample bar chart slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    chart_data = CategoryChartData()
    chart_data.categories = ['A', 'B', 'C']
    chart_data.add_series('Series 1', (10, 20, 30))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5), Inches(6), Inches(4.5), chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM

def add_pie_chart_slide(prs, title):
    """Add a sample pie chart slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    chart_data = CategoryChartData()
    chart_data.categories = ['X', 'Y', 'Z']
    chart_data.add_series('Series 1', (40, 30, 30))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(1), Inches(1.5), Inches(6), Inches(4.5), chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT

def add_text_box_slide(prs, title, text):
    """Add a slide with a text box."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(6), Inches(4))
    textbox.text_frame.text = text

def add_image_slide(prs, title, img_path):
    """Add a slide with an image."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(6), height=Inches(4))

def add_table_slide(prs, title):
    """Add a slide with a table."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    # Create a table with 3 rows and 4 columns
    rows = 3
    cols = 4
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(6)
    height = Inches(3)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    for col in range(cols):
        table.columns[col].width = Inches(1.5)

    # Populate the table with sample data
    data = [
        ["Header 1", "Header 2", "Header 3", "Header 4"],
        ["Row 1, Col 1", "Row 1, Col 2", "Row 1, Col 3", "Row 1, Col 4"],
        ["Row 2, Col 1", "Row 2, Col 2", "Row 2, Col 3", "Row 2, Col 4"]
    ]

    for row in range(rows):
        for col in range(cols):
            table.cell(row, col).text = data[row][col]

def generate_presentation(topic, description, image_file=None):
    """Generate a PowerPoint presentation based on user input."""
    prs = Presentation()
    add_title_slide(prs, topic, "Generated using Streamlit & python-pptx")

    # Generate dynamic content from GlamaAI
    generated_content = generate_api_content(topic, description)

    add_text_box_slide(prs, "Detailed Content from GlamaAI", generated_content)
    add_chart_slide(prs, "Bar Chart Representation")
    add_pie_chart_slide(prs, "Pie Chart Breakdown")
    add_table_slide(prs, "Data Table")

    if image_file:
        add_image_slide(prs, "Uploaded Image", image_file)

    filename = "generated_presentation.pptx"
    prs.save(filename)
    return filename


# Streamlit UI
st.title("📊 PowerPoint Generator")

topic = st.text_input("Enter Topic:")
description = st.text_area("Enter Description:")
uploaded_image = st.file_uploader("Upload an image (optional)", type=["png", "jpg", "jpeg"])

if st.button("Generate PPT"):
    if topic and description:
        img_path = None
        if uploaded_image:
            img_path = f"temp_{uploaded_image.name}"
            with open(img_path, "wb") as f:
                f.write(uploaded_image.getbuffer())

        pptx_file = generate_presentation(topic, description, img_path)

        with open(pptx_file, "rb") as file:
            st.download_button(
                label="📥 Download Presentation",
                data=file,
                file_name=pptx_file,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )

        if img_path:
            os.remove(img_path)  # Cleanup temp file
    else:
        st.warning("⚠️ Please enter both a topic and description.")
